from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT.parent / "submission" / "GapMind_作品方案_v1.pptx"

FONT = "Noto Sans SC"
FONT_BOLD = "Noto Sans SC Bold"

NAVY = "0B1020"
INK = "111A30"
PANEL = "151F38"
PANEL_2 = "1A2745"
LINE = "2D3B61"
WHITE = "F7FAFF"
MUTED = "B7C2D9"
CYAN = "55D6FF"
TEAL = "51D1B6"
PURPLE = "9C8CFF"
ORANGE = "FFB86C"
RED = "FF7A90"
GREEN = "79E2A7"

W = 13.333
H = 7.5


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run(run, size: float, color: str = WHITE, bold: bool = False, font: str | None = None):
    run.font.name = font or (FONT_BOLD if bold else FONT)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    rpr = run._r.get_or_add_rPr()
    rpr.set(qn("a:latin"), font or (FONT_BOLD if bold else FONT))
    rpr.set(qn("a:ea"), font or (FONT_BOLD if bold else FONT))
    rpr.set(qn("a:cs"), font or (FONT_BOLD if bold else FONT))


def add_text(slide, x, y, w, h, text, size=16, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.06,
             font=None, line_spacing=1.05):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraphs = str(text).split("\n")
    for index, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        set_run(run, size, color, bold, font)
    return box


def shape(slide, kind, x, y, w, h, fill=PANEL, line_color=LINE,
          radius=False, transparency=0):
    item = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    item.fill.solid()
    item.fill.fore_color.rgb = rgb(fill)
    if transparency:
        item.fill.transparency = transparency
    item.line.color.rgb = rgb(line_color)
    item.line.width = Pt(0.8)
    return item


def rect(slide, x, y, w, h, fill=PANEL, line_color=LINE, transparency=0):
    return shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line_color, transparency=transparency)


def rounded(slide, x, y, w, h, fill=PANEL, line_color=LINE, transparency=0):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line_color, radius=True, transparency=transparency)


def ellipse(slide, x, y, w, h, fill=PANEL, line_color=LINE):
    return shape(slide, MSO_SHAPE.OVAL, x, y, w, h, fill, line_color)


def line(slide, x1, y1, x2, y2, color=LINE, width=1.5):
    item = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    item.line.color.rgb = rgb(color)
    item.line.width = Pt(width)
    return item


def arrow(slide, x1, y1, x2, y2, color=CYAN, width=1.8):
    item = line(slide, x1, y1, x2, y2, color, width)
    item.line.end_arrowhead = True
    return item


def pill(slide, x, y, w, text, color=CYAN, fill=PANEL_2):
    rounded(slide, x, y, w, 0.32, fill, color)
    add_text(slide, x, y + 0.01, w, 0.28, text, 9, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 0.02)


def dot(slide, x, y, color=CYAN, size=0.12):
    ellipse(slide, x, y, size, size, color, color)


def label_value(slide, x, y, label, value, color=CYAN, w=2.2):
    add_text(slide, x, y, w, 0.22, label.upper(), 8, color, True, margin=0.01)
    add_text(slide, x, y + 0.22, w, 0.32, value, 14, WHITE, True, margin=0.01)


def card(slide, x, y, w, h, title, body, accent=CYAN, body_size=11, title_size=15):
    rounded(slide, x, y, w, h, PANEL, LINE)
    rect(slide, x, y, 0.06, h, accent, accent)
    add_text(slide, x + 0.18, y + 0.16, w - 0.28, 0.34, title, title_size, WHITE, True, margin=0.01)
    add_text(slide, x + 0.18, y + 0.58, w - 0.3, h - 0.68, body, body_size, MUTED, margin=0.01, line_spacing=1.12)


def node(slide, x, y, w, h, title, sub, accent=CYAN, title_size=13):
    rounded(slide, x, y, w, h, PANEL_2, accent)
    dot(slide, x + 0.16, y + 0.18, accent, 0.12)
    add_text(slide, x + 0.36, y + 0.12, w - 0.46, 0.25, title, title_size, WHITE, True, margin=0.01)
    add_text(slide, x + 0.16, y + 0.48, w - 0.3, h - 0.55, sub, 9, MUTED, margin=0.01, line_spacing=1.08)


def set_bg(slide, color=NAVY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def chrome(slide, number, section, dark=False):
    color = "243151" if dark else LINE
    line(slide, 0.55, 7.08, 12.78, 7.08, color, 0.7)
    add_text(slide, 0.58, 7.15, 4.2, 0.18, "GAPMIND  /  作品方案 v1", 7.5, MUTED, True, margin=0.01)
    add_text(slide, 10.1, 7.15, 2.65, 0.18, f"{section}   {number:02d} / 14", 7.5, MUTED, False, PP_ALIGN.RIGHT, margin=0.01)


def title_block(slide, kicker, title, subtitle=None):
    add_text(slide, 0.68, 0.45, 4.8, 0.22, kicker.upper(), 9, CYAN, True, margin=0.01)
    add_text(slide, 0.68, 0.76, 11.9, 0.55, title, 27, WHITE, True, margin=0.01)
    if subtitle:
        add_text(slide, 0.7, 1.42, 11.7, 0.35, subtitle, 11, MUTED, margin=0.01)


def add_slide(prs, number, section, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, INK if dark else NAVY)
    chrome(slide, number, section, dark)
    return slide


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    # Decorative research graph.
    for x1, y1, x2, y2, c in [
        (8.8, 1.15, 11.25, 1.72, CYAN), (8.8, 1.15, 9.35, 3.55, PURPLE),
        (11.25, 1.72, 11.75, 3.95, TEAL), (9.35, 3.55, 11.75, 3.95, ORANGE),
        (9.35, 3.55, 8.15, 5.15, CYAN), (11.75, 3.95, 10.3, 5.45, PURPLE),
        (8.15, 5.15, 10.3, 5.45, TEAL),
    ]:
        line(slide, x1, y1, x2, y2, c, 1.2)
    for x, y, c, s in [(8.6, 0.95, CYAN, 0.38), (11.05, 1.52, PURPLE, 0.28),
                       (9.15, 3.35, TEAL, 0.32), (11.55, 3.72, ORANGE, 0.42),
                       (7.95, 4.95, CYAN, 0.25), (10.1, 5.25, PURPLE, 0.3)]:
        ellipse(slide, x, y, s, s, c, c)
    rounded(slide, 0.72, 0.72, 1.05, 0.36, PANEL_2, CYAN)
    add_text(slide, 0.72, 0.75, 1.05, 0.29, "参赛作品方案", 9, CYAN, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE, 0.01)
    add_text(slide, 0.72, 1.55, 6.8, 0.92, "GapMind", 46, WHITE, True, margin=0.01)
    add_text(slide, 0.76, 2.56, 7.1, 0.75, "面向图机器学习 / 图神经网络的\n证据驱动科研创新助手", 23, CYAN, True, margin=0.01, line_spacing=1.12)
    add_text(slide, 0.78, 3.75, 6.4, 0.55, "从论文证据，到研究机会，再到可审阅的计划与代码草稿。", 14, MUTED, margin=0.01)
    pill(slide, 0.78, 5.2, 1.48, "文献理解", CYAN)
    pill(slide, 2.38, 5.2, 1.65, "机会核验", TEAL)
    pill(slide, 4.15, 5.2, 1.78, "计划与代码", PURPLE)
    add_text(slide, 0.78, 6.35, 6.5, 0.28, "AI 负责候选和草稿，关键研究判断由用户确认。", 10, WHITE, True, margin=0.01)
    add_text(slide, 10.45, 6.65, 2.35, 0.2, "yx_dev  ·  2026.08", 8, MUTED, False, PP_ALIGN.RIGHT, margin=0.01)


def slide_02(prs):
    slide = add_slide(prs, 2, "问题")
    title_block(slide, "01 / WHY", "研究创新的瓶颈，不是信息不够，而是证据链不完整。", "图机器学习研究快速迭代，研究者需要同时处理理解、比较、核验和落地。")
    card(slide, 0.72, 2.15, 3.78, 3.35, "01 读懂与比较", "论文方法、数据集和实验设置分散在不同文献中。\n\n研究者需要反复定位原文、对齐术语，才能回答“到底改进了什么”。", CYAN, 12)
    card(slide, 4.78, 2.15, 3.78, 3.35, "02 找到真空地带", "相似工作可能已经覆盖想法，反证也可能藏在边界条件里。\n\n只凭相似度或一句摘要，容易把“未知”误判成“新颖”。", ORANGE, 12)
    card(slide, 8.84, 2.15, 3.78, 3.35, "03 从想法走向实验", "研究想法还需要变成可证伪问题、证据要求、实验计划和代码草稿。\n\n普通 Chat 往往给出漂亮但不可追溯的答案。", PURPLE, 12)
    add_text(slide, 0.78, 6.08, 11.8, 0.52, "普通 Chat 解决“生成一段文字”，论文管理工具解决“保存一批文件”；GapMind 连接的是“证据 → 判断 → 行动”。", 16, WHITE, True, margin=0.01)


def slide_03(prs):
    slide = add_slide(prs, 3, "定位")
    title_block(slide, "02 / POSITIONING", "把科研辅助从“问答工具”推进到“可审阅的研究协作”。")
    rounded(slide, 0.72, 2.0, 11.88, 1.25, PANEL_2, CYAN)
    add_text(slide, 0.98, 2.27, 11.35, 0.7, "GapMind 是面向计算机科学—图机器学习/图神经网络研究场景的证据驱动科研创新助手，帮助研究生、科研助理和导师完成文献理解、研究机会核验、可证伪研究计划设计和代码草稿生成。", 16, WHITE, True, margin=0.01, line_spacing=1.15)
    add_text(slide, 0.78, 3.75, 2.2, 0.25, "两类核心用户", 10, CYAN, True, margin=0.01)
    card(slide, 0.72, 4.1, 5.6, 1.65, "研究生 / 科研助理", "快速建立领域地图；核对方法与反证；把想法变成下一步可执行的研究任务。", TEAL, 12)
    card(slide, 6.72, 4.1, 5.88, 1.65, "导师 / 教师", "审阅候选的新颖性边界；追踪证据质量；把讨论沉淀为可复用的研究计划与代码草稿。", PURPLE, 12)
    pill(slide, 0.8, 6.18, 1.7, "证据问答", CYAN)
    pill(slide, 2.68, 6.18, 2.0, "机会核验", TEAL)
    pill(slide, 4.86, 6.18, 2.45, "计划 / 代码草稿", PURPLE)


def slide_04(prs):
    slide = add_slide(prs, 4, "闭环")
    title_block(slide, "03 / RESEARCH LOOP", "一条从论文证据到研究行动的闭环。", "每一步都保留来源、状态和人工确认点，避免把生成结果直接当作研究结论。")
    steps = [
        ("导入", "论文 / 资料", CYAN), ("理解", "解析 · 抽取 · 索引", CYAN),
        ("问答", "引用原文证据", TEAL), ("检索", "支持 · 相似 · 反证", TEAL),
        ("核验", "外部新颖性", ORANGE), ("候选", "Opportunity", ORANGE),
        ("收窄", "Critic", PURPLE), ("门控", "Evidence Gate", PURPLE),
        ("确认", "HITL", RED), ("行动", "计划 · 代码草稿", GREEN),
    ]
    x0, y0, nw, nh, gap = 0.72, 2.25, 2.26, 0.86, 0.23
    for i, (top, sub, c) in enumerate(steps):
        row = 0 if i < 5 else 1
        col = i if i < 5 else i - 5
        x = x0 + col * (nw + gap)
        y = y0 + row * 1.52
        rounded(slide, x, y, nw, nh, PANEL_2, c)
        add_text(slide, x + 0.14, y + 0.13, nw - 0.28, 0.24, top, 13, c, True, margin=0.01)
        add_text(slide, x + 0.14, y + 0.45, nw - 0.28, 0.22, sub, 9.5, WHITE, margin=0.01)
        if i < 4:
            arrow(slide, x + nw, y + nh / 2, x + nw + gap - 0.06, y + nh / 2, LINE, 1.4)
        if i == 4:
            arrow(slide, x + nw / 2, y + nh, x + nw / 2, y + 1.52 - 0.15, LINE, 1.4)
        if i >= 5 and i < 9:
            arrow(slide, x + nw, y + nh / 2, x + nw + gap - 0.06, y + nh / 2, LINE, 1.4)
    rounded(slide, 0.72, 5.62, 5.75, 0.82, PANEL, LINE)
    add_text(slide, 0.92, 5.83, 5.35, 0.2, "AI 负责", 10, CYAN, True, margin=0.01)
    add_text(slide, 1.7, 5.78, 4.45, 0.3, "检索 · 归纳 · 候选 · 草稿 · 缺口提示", 12, WHITE, margin=0.01)
    rounded(slide, 6.72, 5.62, 5.88, 0.82, PANEL, LINE)
    add_text(slide, 6.92, 5.83, 5.35, 0.2, "人负责", 10, RED, True, margin=0.01)
    add_text(slide, 7.7, 5.78, 4.55, 0.3, "证据判断 · 候选取舍 · 计划确认 · 实验决策", 12, WHITE, margin=0.01)


def slide_05(prs):
    slide = add_slide(prs, 5, "能力")
    title_block(slide, "04 / CAPABILITIES", "四个能力模块，共用一套证据与状态语言。")
    card(slide, 0.72, 2.05, 5.72, 1.82, "Workspace Chat", "多轮问答、来源隔离、引用质量检查；当检索不到足够证据时，系统明确收敛，不用猜测补全。", CYAN, 11.5)
    card(slide, 6.72, 2.05, 5.88, 1.82, "Discover", "Planner → Evidence → External → Opportunity → Critic → Gate，把研究机会拆成可观察的步骤。", TEAL, 11.5)
    card(slide, 0.72, 4.2, 5.72, 1.82, "Evidence Passport", "记录独立论文、全文 / 元数据、支持 / 反证、门控版本和人工状态，让“为什么相信”可以回溯。", ORANGE, 11.5)
    card(slide, 6.72, 4.2, 5.88, 1.82, "Agent 产物", "研究计划、代码草稿以及后续分析 / 写作 / 审稿回复均保留 AI 生成状态，支持预览、下载和人工接管。", PURPLE, 11.5)
    add_text(slide, 0.78, 6.48, 11.7, 0.22, "当前产品重点：证据驱动的科研创新闭环；GraphRAG、多模态、代码执行和规模化多租户属于后续验证方向。", 10, MUTED, margin=0.01)


def slide_06(prs):
    slide = add_slide(prs, 6, "架构")
    title_block(slide, "05 / ARCHITECTURE", "垂类抽取、可靠检索和受控生成，组成可解释的技术底座。")
    # left source / middle services / right output
    node(slide, 0.72, 2.2, 2.0, 0.82, "论文与资料", "PDF · workspace", CYAN)
    node(slide, 0.72, 3.25, 2.0, 0.82, "垂类抽取", "实体 · 关系 · 证据", TEAL)
    node(slide, 0.72, 4.3, 2.0, 0.82, "外部知识", "Semantic Scholar", ORANGE)
    rounded(slide, 3.25, 1.95, 4.7, 3.75, PANEL, LINE)
    add_text(slide, 3.5, 2.2, 4.2, 0.25, "科研工作空间服务层", 14, WHITE, True, margin=0.01)
    for x, y, w, text, c in [(3.5, 2.82, 1.85, "Chat / RAG", CYAN), (5.55, 2.82, 2.05, "Discover", TEAL),
                              (3.5, 3.8, 1.85, "Evidence", ORANGE), (5.55, 3.8, 2.05, "Agent / HITL", PURPLE),
                              (3.5, 4.78, 1.85, "API / Auth", RED), (5.55, 4.78, 2.05, "Timeline", GREEN)]:
        rounded(slide, x, y, w, 0.58, PANEL_2, c)
        add_text(slide, x, y + 0.13, w, 0.25, text, 11, c, True, PP_ALIGN.CENTER, margin=0.01)
    node(slide, 8.5, 2.2, 1.95, 0.82, "PostgreSQL", "业务状态 / 证据", CYAN)
    node(slide, 8.5, 3.25, 1.95, 0.82, "Milvus", "dense retrieval", TEAL)
    node(slide, 8.5, 4.3, 1.95, 0.82, "Redis / Celery", "异步任务", ORANGE)
    node(slide, 10.92, 2.2, 1.7, 0.82, "LLM", "受控生成", PURPLE)
    node(slide, 10.92, 3.25, 1.7, 0.82, "Reranker", "二阶段排序", CYAN)
    node(slide, 10.92, 4.3, 1.7, 0.82, "React", "可视化交互", GREEN)
    for y in [2.61, 3.66, 4.71]:
        arrow(slide, 2.72, y, 3.25, y, LINE, 1.2)
        arrow(slide, 7.95, y, 8.5, y, LINE, 1.2)
    for y in [2.61, 3.66, 4.71]:
        arrow(slide, 10.45, y, 10.92, y, LINE, 1.2)
    rounded(slide, 0.72, 6.13, 11.9, 0.6, PANEL_2, LINE)
    add_text(slide, 0.93, 6.31, 11.45, 0.22, "工程约束：Chat 当前生产链路为 dense retrieval + reranker；来源严格区分 [En] / [Pn] / [Dn] / [Cn]；LLM 调用统一 disable_thinking=True。", 10, WHITE, margin=0.01)


def slide_07(prs):
    slide = add_slide(prs, 7, "门控")
    title_block(slide, "06 / AGENT + GATE", "多智能体的价值，在于分工、挑战与收窄，而不是堆叠数量。")
    # Main state path
    xs = [0.75, 2.78, 4.81, 6.84, 8.87, 10.9]
    labels = [("Planner", "拆研究轴", CYAN), ("Evidence", "找支持 / 反证", TEAL),
              ("External", "实时 / 缓存 / 失败", ORANGE), ("Opportunity", "形成候选", PURPLE),
              ("Critic", "keep / narrow / reject", RED), ("Gate", "独立全文 + 覆盖", GREEN)]
    for i, (t, s, c) in enumerate(labels):
        rounded(slide, xs[i], 2.25, 1.63, 0.98, PANEL_2, c)
        add_text(slide, xs[i], 2.42, 1.63, 0.22, t, 11, c, True, PP_ALIGN.CENTER, margin=0.01)
        add_text(slide, xs[i] + 0.08, 2.76, 1.47, 0.24, s, 8.5, WHITE, False, PP_ALIGN.CENTER, margin=0.01)
        if i < len(labels) - 1:
            arrow(slide, xs[i] + 1.63, 2.74, xs[i + 1] - 0.08, 2.74, LINE, 1.2)
    # branches
    rounded(slide, 3.25, 3.9, 2.55, 1.15, PANEL, ORANGE)
    add_text(slide, 3.45, 4.1, 2.15, 0.22, "外部服务部分成功", 11, ORANGE, True, margin=0.01)
    add_text(slide, 3.45, 4.42, 2.15, 0.4, "标识状态，保留已获得证据\n不伪装成完整核验", 9.5, MUTED, margin=0.01)
    rounded(slide, 6.28, 3.9, 2.55, 1.15, PANEL, RED)
    add_text(slide, 6.48, 4.1, 2.15, 0.22, "Critic = narrow", 11, RED, True, margin=0.01)
    add_text(slide, 6.48, 4.42, 2.15, 0.4, "收窄研究轴、补查证据\n或拒绝不成立的候选", 9.5, MUTED, margin=0.01)
    rounded(slide, 9.32, 3.9, 3.3, 1.15, PANEL, GREEN)
    add_text(slide, 9.52, 4.1, 2.9, 0.22, "Gate → HITL", 11, GREEN, True, margin=0.01)
    add_text(slide, 9.52, 4.42, 2.9, 0.4, "AI 给出证据护照和缺口\n用户决定是否进入研究计划", 9.5, MUTED, margin=0.01)
    arrow(slide, 4.65, 3.23, 4.45, 3.9, ORANGE, 1.2)
    arrow(slide, 7.65, 3.23, 7.55, 3.9, RED, 1.2)
    arrow(slide, 10.05, 3.23, 10.8, 3.9, GREEN, 1.2)
    add_text(slide, 0.78, 5.7, 11.7, 0.5, "关键设计：失败有分类，候选有状态，进入长期资产前必须通过证据门和人工确认。", 16, WHITE, True, margin=0.01)


def slide_08(prs):
    slide = add_slide(prs, 8, "可信")
    title_block(slide, "07 / TRUST", "让每一个结论，都能回到来源、版本和确认状态。")
    # passport chain
    stages = [("回答", "引用 [En]", CYAN), ("原文", "论文 / 证据片段", TEAL),
              ("候选", "Opportunity", ORANGE), ("护照", "Evidence Passport", PURPLE),
              ("计划", "输入版本 + 人工确认", GREEN)]
    for i, (t, s, c) in enumerate(stages):
        x = 0.78 + i * 2.45
        ellipse(slide, x, 2.32, 0.58, 0.58, c, c)
        add_text(slide, x, 2.47, 0.58, 0.2, str(i + 1), 12, NAVY, True, PP_ALIGN.CENTER, margin=0.01)
        add_text(slide, x - 0.08, 3.08, 1.0, 0.22, t, 12, WHITE, True, PP_ALIGN.CENTER, margin=0.01)
        add_text(slide, x - 0.52, 3.42, 1.9, 0.38, s, 9.5, MUTED, False, PP_ALIGN.CENTER, margin=0.01)
        if i < len(stages) - 1:
            arrow(slide, x + 0.64, 2.61, x + 2.36, 2.61, LINE, 1.6)
    card(slide, 0.72, 4.35, 3.78, 1.45, "生成边界", "AI 生成内容显式标识；未确认产物不自动成为长期事实。", CYAN, 11)
    card(slide, 4.78, 4.35, 3.78, 1.45, "安全边界", "敏感材料提示脱敏与远程调用风险；代码默认只预览 / 下载。", RED, 11)
    card(slide, 8.84, 4.35, 3.78, 1.45, "诚实边界", "不伪造引用、实验结果和用户结论；静态检查不等于真实执行。", ORANGE, 11)
    add_text(slide, 0.78, 6.28, 11.7, 0.32, "Evidence Passport 不是一个“置信度分数”，而是一份可审阅的证据清单：覆盖什么、缺少什么、谁确认过。", 12, WHITE, True, margin=0.01)


def slide_09(prs):
    slide = add_slide(prs, 9, "验证")
    title_block(slide, "08 / QUALITY", "质量证明从典型问题开始：回答正确，更要证据对得上。", "本页预留冻结 Gold 与人工评审结果；当前版本不填未经验证的效果数字。")
    headers = ["问题类型", "验证重点", "结果呈现"]
    widths = [2.25, 5.25, 4.15]
    x = 0.72
    for h, w in zip(headers, widths):
        rounded(slide, x, 2.1, w, 0.55, PANEL_2, CYAN)
        add_text(slide, x + 0.12, 2.25, w - 0.24, 0.22, h, 11, CYAN, True, margin=0.01)
        x += w + 0.12
    rows = [
        ("概念 / 方法规范题", "术语、方法机制、适用边界", "标准答案摘要 + [En] 原文 + 人工结论"),
        ("跨论文比较题", "比较维度是否完整，是否混淆来源", "多篇来源 + 关键差异 + 引用一致性"),
        ("证据不足题", "是否识别缺口，避免绝对化结论", "明确“不足以判断” + 缺口说明"),
    ]
    y = 2.8
    for i, row in enumerate(rows):
        c = [CYAN, TEAL, ORANGE][i]
        x = 0.72
        for j, (value, w) in enumerate(zip(row, widths)):
            rounded(slide, x, y, w, 0.84, PANEL, LINE)
            add_text(slide, x + 0.12, y + 0.16, w - 0.24, 0.5, value, 10.5 if j else 11, WHITE if j else c, j == 0, margin=0.01, line_spacing=1.1)
            x += w + 0.12
        y += 1.02
    rounded(slide, 0.72, 5.98, 11.9, 0.66, PANEL_2, LINE)
    add_text(slide, 0.94, 6.18, 11.45, 0.25, "待填：Gold 版本 / 评审人数 / 支持率 / 引用有效率 / 失败案例。所有数字必须回到冻结 Gold、人工评审或用户记录。", 10, ORANGE, True, margin=0.01)


def slide_10(prs):
    slide = add_slide(prs, 10, "用户")
    title_block(slide, "09 / USER VALIDATION", "真实用户效果，是下一阶段必须补齐的证据。", "先记录任务过程和前后对比，再总结效率、质量和继续使用意愿；空缺不以推测数字补齐。")
    for x, title, accent in [(0.72, "用户 A  ·  待填", CYAN), (6.72, "用户 B  ·  待填", TEAL)]:
        rounded(slide, x, 2.15, 5.7, 3.25, PANEL, accent)
        add_text(slide, x + 0.22, 2.38, 5.2, 0.3, title, 15, accent, True, margin=0.01)
        fields = ["身份 / 学科", "任务与场景", "使用次数 / 时长", "前后对比", "反馈 / 授权"]
        for i, field in enumerate(fields):
            yy = 2.88 + i * 0.47
            add_text(slide, x + 0.24, yy, 1.5, 0.2, field, 9.5, MUTED, True, margin=0.01)
            line(slide, x + 1.8, yy + 0.18, x + 5.38, yy + 0.18, LINE, 0.8)
            add_text(slide, x + 1.86, yy - 0.01, 3.4, 0.22, "待填 / 未验证", 9.5, ORANGE, False, margin=0.01)
    add_text(slide, 0.78, 5.85, 11.7, 0.25, "建议指标：任务完成时间、证据核查时间、人工修改次数、发现的风险 / 反证、继续使用意愿。", 11, WHITE, True, margin=0.01)
    rounded(slide, 0.72, 6.25, 11.9, 0.45, PANEL_2, LINE)
    add_text(slide, 0.94, 6.37, 11.4, 0.2, "当前状态：验证框架已建立，真实用户记录与授权待补。", 9.5, ORANGE, True, margin=0.01)


def slide_11(prs):
    slide = add_slide(prs, 11, "创新")
    title_block(slide, "10 / INNOVATION", "创新点不是“模型更大”，而是把研究判断变成可追溯、可协作的过程。")
    innovations = [
        ("证据护照", "把独立全文、支持 / 反证、版本和人工状态放进同一条链路。", CYAN),
        ("Critic + Gate", "研究机会先被挑战和收窄，再决定是否进入长期研究资产。", TEAL),
        ("候选与事实分离", "AI 可以提出候选，但关键结论必须由用户确认。", ORANGE),
        ("失败可解释", "区分失败、部分成功和证据不足，让用户知道下一步怎么做。", PURPLE),
        ("发现到行动", "从研究空白连续推进到可证伪计划和代码草稿。", GREEN),
    ]
    for i, (t, b, c) in enumerate(innovations):
        y = 2.04 + i * 0.84
        rounded(slide, 0.72, y, 2.25, 0.62, PANEL_2, c)
        add_text(slide, 0.88, y + 0.17, 1.93, 0.22, t, 11, c, True, margin=0.01)
        rounded(slide, 3.25, y, 9.36, 0.62, PANEL, LINE)
        add_text(slide, 3.5, y + 0.16, 8.9, 0.25, b, 11.5, WHITE, margin=0.01)
    rounded(slide, 0.72, 6.35, 11.9, 0.46, PANEL_2, LINE)
    add_text(slide, 0.94, 6.48, 11.4, 0.2, "谨慎表述：GraphRAG、混合检索、多模态、代码沙箱等放在后续路线，除非有独立 A/B 和安全证据。", 9.5, MUTED, margin=0.01)


def slide_12(prs):
    slide = add_slide(prs, 12, "推广")
    title_block(slide, "11 / SCALE", "从单实验室试点开始，把垂类科研工作流复制到更多学科。")
    card(slide, 0.72, 2.05, 3.78, 3.45, "首个试点", "学院 / 实验室 / 研究生课程组\n\n以一组受控论文集、一个 Gold 集和少量真实用户开始，先证明过程价值。", CYAN, 12)
    card(slide, 4.78, 2.05, 3.78, 3.45, "复制方式", "替换知识源、Gold、提示模板和垂类抽取器。\n\n工作空间、证据护照、HITL 和 Agent 协议保持复用。", TEAL, 12)
    card(slide, 8.84, 2.05, 3.78, 3.45, "部署与成本", "私有化或受控云；成本来自模型、向量库、外部检索、存储和人工审核。\n\n当前定位：单实验室 / 院系试点。", PURPLE, 12)
    rounded(slide, 0.72, 5.88, 11.9, 0.75, PANEL_2, ORANGE)
    add_text(slide, 0.95, 6.08, 1.05, 0.22, "风险", 10, ORANGE, True, margin=0.01)
    add_text(slide, 1.75, 6.04, 10.45, 0.3, "数据授权 · 学术诚信 · 用户权限 · 模型供应商依赖 · 规模化后的审阅成本", 11, WHITE, margin=0.01)


def slide_13(prs):
    slide = add_slide(prs, 13, "路线图")
    title_block(slide, "12 / ROADMAP", "先把证据和用户验证做实，再扩展更复杂的智能能力。")
    columns = [(0.72, 3.78, "已有工程基础", CYAN, "论文 / 知识 / 检索 / Chat\nDiscover / HITL / 计划 / 代码草稿\n\n后端 481 tests · 前端 56 tests\nTypeScript 与构建通过"),
               (4.78, 3.78, "当前验证中", ORANGE, "两名真实用户效果\n冻结 Gold 扩展与人工评审\n真实外部服务与 Demo 稳定性\n权威知识源治理"),
               (8.84, 3.78, "后续方向", PURPLE, "混合检索 / GraphRAG A/B\n最小多模态或专业数据分析\n受控代码执行\n多租户与规模化")]
    for x, w, t, c, b in columns:
        rounded(slide, x, 2.15, w, 3.65, PANEL, c)
        add_text(slide, x + 0.22, 2.42, w - 0.44, 0.3, t, 15, c, True, margin=0.01)
        add_text(slide, x + 0.22, 3.05, w - 0.44, 2.25, b, 11.5, WHITE, margin=0.01, line_spacing=1.25)
    rounded(slide, 0.72, 6.15, 11.9, 0.55, PANEL_2, LINE)
    add_text(slide, 0.95, 6.32, 11.35, 0.2, "交付原则：每个新增能力，都要同时补齐来源、失败状态、人工边界和可重复验证。", 10.5, WHITE, True, margin=0.01)


def slide_14(prs):
    slide = add_slide(prs, 14, "声明", dark=True)
    title_block(slide, "13 / DECLARATION", "一个人开发，也可以把边界、证据和责任说清楚。", "GapMind 的目标不是替研究者做决定，而是让研究者更快地看见证据、缺口和下一步。")
    rounded(slide, 0.72, 2.25, 7.18, 3.2, PANEL, CYAN)
    add_text(slide, 1.02, 2.58, 6.55, 0.28, "开发与责任边界", 15, CYAN, True, margin=0.01)
    add_text(slide, 1.02, 3.1, 6.35, 1.85, "单人开发，AI 工具辅助分析、起草、编码建议和测试。\n\n用户本人负责最终产品、数据、事实、用户与合规判断。\n\n所有 AI 生成内容按规定标识；关键资产需要人工确认。", 13, WHITE, margin=0.01, line_spacing=1.25)
    rounded(slide, 8.32, 2.25, 4.3, 3.2, PANEL_2, TEAL)
    add_text(slide, 8.62, 2.58, 3.7, 0.28, "下一步", 15, TEAL, True, margin=0.01)
    add_text(slide, 8.62, 3.1, 3.55, 1.85, "补齐真实用户记录\n冻结质量 Gold\n完成真实 Demo 录像\n形成可复核的交付包", 13, WHITE, margin=0.01, line_spacing=1.35)
    add_text(slide, 0.78, 6.25, 11.7, 0.4, "GapMind  ·  让科研创新从“我觉得”走向“我能追溯、我能验证、我能继续做”。", 15, WHITE, True, PP_ALIGN.CENTER, margin=0.01)


def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide_cover(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)
    slide_13(prs)
    slide_14(prs)
    prs.core_properties.title = "GapMind 作品方案 v1"
    prs.core_properties.subject = "面向图机器学习 / 图神经网络的证据驱动科研创新助手"
    prs.core_properties.author = "GapMind"
    prs.core_properties.comments = "初稿：真实效果数据、用户记录和最终演示材料待补。"
    prs.save(OUT)
    print(f"saved: {OUT}")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
